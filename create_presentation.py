from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    title_shape.text = title
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
    return slide

def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    return slide

def add_content_slide(prs, title, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
    
    return slide

def add_example_slide(prs, title, example_text, explanation_list):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    p = tf.add_paragraph()
    p.text = "Example: " + example_text
    p.level = 0
    
    for item in explanation_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
    
    return slide

def add_comparison_slide(prs, title, left_title, right_title, left_items, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Left column
    left = slide.placeholders[1]
    tf = left.text_frame
    tf.text = left_title
    
    for item in left_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
    
    # Right column
    right = slide.placeholders[2]
    tf = right.text_frame
    tf.text = right_title
    
    for item in right_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
    
    return slide

def create_data_representation_presentation():
    prs = Presentation()
    
    # Title slide
    add_title_slide(prs, 
                   "Text Data Representation for Recommendation Systems",
                   "Understanding how to represent text data for effective recommendations")
    
    # Introduction slide
    add_content_slide(prs, 
                     "Introduction",
                     ["Data representation is crucial for recommendation systems",
                      "How we represent text determines what patterns our system can learn",
                      "Different representation techniques have different strengths and weaknesses",
                      "We'll explore: Bag of Words, TF-IDF, and Word Embeddings"])
    
    # Bag of Words slide
    add_content_slide(prs,
                     "Bag of Words (BoW)",
                     ["Represents documents as vectors of word counts",
                      "Ignores word order and context",
                      "Each dimension corresponds to a word in the vocabulary",
                      "Simple to implement but creates sparse, high-dimensional vectors",
                      "Doesn't capture semantic meaning between words"])
    
    # BoW Example slide
    add_example_slide(prs,
                     "Bag of Words - Example",
                     "Document: 'I love machine learning and AI'",
                     ["Vocabulary: {'I', 'love', 'machine', 'learning', 'and', 'AI'}",
                      "Vector representation: [1, 1, 1, 1, 1, 1]",
                      "For document: 'AI and machine learning are transforming industries'",
                      "Vector (same vocabulary): [0, 0, 1, 1, 1, 1, 0, 0] (with 'are', 'transforming', 'industries' added)",
                      "Notice how word order is lost and vectors become sparse as vocabulary grows"])
    
    # TF-IDF slide
    add_content_slide(prs,
                     "TF-IDF (Term Frequency-Inverse Document Frequency)",
                     ["Improves on BoW by weighting terms based on importance",
                      "Term Frequency (TF): How often a word appears in a document",
                      "Inverse Document Frequency (IDF): How rare a word is across all documents",
                      "TF-IDF = TF × IDF",
                      "Reduces impact of common words and emphasizes distinctive terms"])
    
    # TF-IDF Example slide
    add_example_slide(prs,
                     "TF-IDF - Example",
                     "Document 1: 'I love machine learning and AI'",
                     ["Term 'machine' appears in 2 out of 5 documents in our corpus",
                      "TF(machine, doc1) = 1/6 (appears once in a document with 6 words)",
                      "IDF(machine) = log(5/2) ≈ 0.4 (appears in 2 out of 5 documents)",
                      "TF-IDF(machine, doc1) = 1/6 × 0.4 ≈ 0.067",
                      "Common words like 'and' will have lower TF-IDF scores than distinctive terms like 'AI'"])
    
    # Word Embeddings slide
    add_content_slide(prs,
                     "Word Embeddings",
                     ["Represent words as dense vectors in a continuous vector space",
                      "Words with similar meanings are positioned close to each other",
                      "Vectors capture semantic relationships between words",
                      "Pre-trained models: Word2Vec, GloVe, BERT, Sentence Transformers",
                      "Captures semantic meaning but requires more computational resources"])
    
    # Word Embeddings Example slide
    add_example_slide(prs,
                     "Word Embeddings - Example",
                     "Using pre-trained embeddings to represent 'machine learning'",
                     ["Word 'machine' might be represented as: [0.2, -0.6, 0.1, 0.3, ...]",
                      "Word 'learning' might be represented as: [0.3, -0.4, 0.2, 0.1, ...]",
                      "Document embedding could be the average: [0.25, -0.5, 0.15, 0.2, ...]",
                      "Semantic relationships: vector('king') - vector('man') + vector('woman') ≈ vector('queen')",
                      "Similar concepts like 'AI' and 'machine learning' will have similar vectors"])
    
    # Sentence Transformers Example slide
    add_example_slide(prs,
                     "Sentence Transformers - Example",
                     "Encoding entire sentences with context awareness",
                     ["Input: 'Natural language processing is a subfield of AI'",
                      "Output: Dense vector of typically 384-768 dimensions",
                      "Captures contextual meaning of words in relation to each other",
                      "Example vector (truncated): [0.12, -0.34, 0.56, 0.78, -0.91, ...]",
                      "Similar sentences will have high cosine similarity even with different words"])
    
    # Comparison slide
    add_comparison_slide(prs,
                        "Comparing Representation Techniques",
                        "Traditional Methods (BoW, TF-IDF)",
                        "Modern Methods (Word Embeddings)",
                        ["Simple to implement",
                         "Computationally efficient",
                         "Works well for basic tasks",
                         "Sparse, high-dimensional vectors",
                         "No semantic understanding"],
                        ["Captures semantic meaning",
                         "Dense, lower-dimensional vectors",
                         "Represents word relationships",
                         "Better performance on complex tasks",
                         "Requires more computational resources"])
    
    # Document Similarity slide
    add_content_slide(prs,
                     "Document Similarity Measures",
                     ["Cosine Similarity: Measures the angle between vectors (most common)",
                      "Euclidean Distance: Straight-line distance between points",
                      "Jaccard Similarity: Ratio of intersection to union of sets",
                      "Manhattan Distance: Sum of absolute differences",
                      "Similarity calculation is key for finding relevant recommendations"])
    
    # Similarity Example slide
    add_example_slide(prs,
                     "Document Similarity - Example",
                     "Comparing two documents using cosine similarity",
                     ["Document 1: 'I love machine learning and AI'",
                      "Document 2: 'AI and machine learning are transforming industries'",
                      "TF-IDF vectors: [0.1, 0.3, 0.2, 0.4, 0.1, 0.5] and [0.5, 0.1, 0.2, 0.4, 0, 0]",
                      "Cosine similarity = (0.1×0.5 + 0.3×0.1 + 0.2×0.2 + 0.4×0.4 + ...) / (||vec1|| × ||vec2||) ≈ 0.75",
                      "Higher similarity (0.75) indicates documents are related despite different words"])
    
    # Building a Recommendation System slide
    add_content_slide(prs,
                     "Building a Simple Recommendation System",
                     ["1. Represent all documents as vectors",
                      "2. For a given query, find the most similar documents",
                      "3. Recommend the top N most similar documents",
                      "4. Evaluate and refine the system",
                      "5. Consider hybrid approaches for better performance"])
    
    # Recommendation Example slide
    add_example_slide(prs,
                     "Recommendation System - Example",
                     "Content-based recommendation for articles",
                     ["User reads article: 'Introduction to Machine Learning'",
                      "System represents this article as a vector using TF-IDF or embeddings",
                      "Calculates similarity with all other articles in the database",
                      "Returns top 3 most similar articles: 'Deep Learning Basics', 'AI Fundamentals', 'Neural Networks Explained'",
                      "System improves as user provides feedback on recommendations"])
    
    # Practical Implementation slide
    add_content_slide(prs,
                     "Practical Implementation",
                     ["Text preprocessing: lowercase, remove punctuation, stopwords",
                      "Feature extraction: BoW, TF-IDF using scikit-learn",
                      "Word embeddings: Using Sentence Transformers",
                      "Similarity calculation: Using cosine similarity",
                      "Recommendation function: Return top N similar items"])
    
    # Conclusion slide
    add_content_slide(prs,
                     "Conclusion",
                     ["Data representation is fundamental to recommendation systems",
                      "Different techniques offer different trade-offs",
                      "Modern embedding techniques generally outperform traditional methods",
                      "Preprocessing and similarity measures are also important",
                      "Experiment with different approaches for your specific use case"])
    
    # Save the presentation
    prs.save("c:/Users/dario/projects/ai-recommender/Data_Representation_for_Recommendations.pptx")
    print("Presentation created successfully at: c:/Users/dario/projects/ai-recommender/Data_Representation_for_Recommendations.pptx")

if __name__ == "__main__":
    create_data_representation_presentation()