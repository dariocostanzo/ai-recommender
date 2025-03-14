from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    title_shape.text = title
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
    return slide

def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    return slide

def add_content_slide(prs, title, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
    
    return slide

def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
    
    return slide

def add_comparison_slide(prs, title, left_title, right_title, left_items, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Left column
    left = slide.placeholders[1]
    tf = left.text_frame
    tf.text = left_title
    
    for item in left_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
    
    # Right column
    right = slide.placeholders[2]
    tf = right.text_frame
    tf.text = right_title
    
    for item in right_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
    
    return slide

def create_data_representation_presentation():
    prs = Presentation()
    
    # Title slide
    add_title_slide(prs, 
                   "Text Data Representation for Recommendation Systems",
                   "Understanding how to represent text data for effective recommendations")
    
    # Introduction slide
    add_content_slide(prs, 
                     "Introduction",
                     ["Data representation is crucial for recommendation systems",
                      "How we represent text determines what patterns our system can learn",
                      "Different representation techniques have different strengths and weaknesses",
                      "We'll explore: Bag of Words, TF-IDF, and Word Embeddings"])
    
    # Bag of Words slide
    add_bullet_slide(prs,
                     "Bag of Words (BoW)",
                     ["• Each document is represented as a vector of word counts",
                      "• The position and order of words are ignored",
                      "• Each dimension corresponds to a word in the vocabulary",
                      "",
                      "Advantages:",
                      "• Simple to understand and implement",
                      "• Works well for basic text classification",
                      "",
                      "Disadvantages:",
                      "• Ignores word order and context",
                      "• Results in sparse, high-dimensional vectors",
                      "• Doesn't capture semantic meaning"])
    
    # TF-IDF slide
    add_bullet_slide(prs,
                     "TF-IDF (Term Frequency-Inverse Document Frequency)",
                     ["• Improves on BoW by weighting terms based on importance",
                      "• Term Frequency (TF): How often a word appears in a document",
                      "• Inverse Document Frequency (IDF): How rare a word is across all documents",
                      "• TF-IDF = TF × IDF",
                      "",
                      "Advantages:",
                      "• Reduces impact of common words",
                      "• Gives more weight to distinctive terms",
                      "• Better performance than simple BoW",
                      "",
                      "Disadvantages:",
                      "• Still uses sparse vectors",
                      "• Doesn't capture semantic relationships between words",
                      "• No understanding of word context"])
    
    # Word Embeddings slide
    add_bullet_slide(prs,
                     "Word Embeddings",
                     ["• Words represented as dense vectors in a continuous vector space",
                      "• Words with similar meanings are positioned close to each other",
                      "• Vectors capture semantic relationships",
                      "• Pre-trained models: Word2Vec, GloVe, BERT, Sentence Transformers",
                      "",
                      "Advantages:",
                      "• Captures semantic meaning",
                      "• Lower-dimensional representation",
                      "• Can represent relationships between words (e.g., king - man + woman ≈ queen)",
                      "",
                      "Disadvantages:",
                      "• More complex to implement",
                      "• Requires more computational resources",
                      "• May need large amounts of training data"])
    
    # Document Similarity slide
    add_bullet_slide(prs,
                     "Document Similarity Measures",
                     ["Cosine Similarity:",
                      "• Measures the cosine of the angle between two vectors",
                      "• Ranges from -1 (completely opposite) to 1 (exactly the same)",
                      "• Formula: cos(θ) = (A·B) / (||A|| × ||B||)",
                      "",
                      "Euclidean Distance:",
                      "• The straight-line distance between two points",
                      "• Smaller values indicate more similar documents",
                      "• Formula: d(A, B) = √(Σ(Ai - Bi)²)"])
    
    # Comparison slide
    add_comparison_slide(prs,
                        "Comparing Representation Techniques",
                        "Traditional Methods (BoW, TF-IDF)",
                        "Modern Methods (Word Embeddings)",
                        ["Simple to implement",
                         "Computationally efficient",
                         "Works well for basic tasks",
                         "Sparse, high-dimensional vectors",
                         "No semantic understanding"],
                        ["Captures semantic meaning",
                         "Dense, lower-dimensional vectors",
                         "Represents word relationships",
                         "Better performance on complex tasks",
                         "Requires more computational resources"])
    
    # Building a Recommendation System slide
    add_bullet_slide(prs,
                     "Building a Simple Recommendation System",
                     ["1. Represent all documents as vectors",
                      "2. For a given query or document, find the most similar documents",
                      "3. Recommend the top N most similar documents",
                      "4. Evaluate and refine the system",
                      "5. Consider hybrid approaches for better performance"])
    
    # Practical Implementation slide
    add_bullet_slide(prs,
                     "Practical Implementation",
                     ["Text preprocessing:",
                      "• Lowercase conversion, punctuation removal",
                      "• Stopword removal, stemming, lemmatization",
                      "",
                      "Feature extraction:",
                      "• BoW, TF-IDF using scikit-learn",
                      "• Word embeddings using Sentence Transformers",
                      "",
                      "Similarity calculation:",
                      "• Using cosine similarity to find similar documents",
                      "",
                      "Recommendation function:",
                      "• Return top N similar items based on similarity scores"])
    
    # Conclusion slide
    add_bullet_slide(prs,
                     "Conclusion",
                     ["• Data representation is fundamental to recommendation systems",
                      "• Different techniques offer different trade-offs",
                      "• Modern embedding techniques generally outperform traditional methods",
                      "• Preprocessing and similarity measures are also important",
                      "• Experiment with different approaches for your specific use case",
                      "",
                      "Next steps:",
                      "• Explore more advanced embedding techniques",
                      "• Implement hybrid recommendation approaches",
                      "• Evaluate system performance with real-world data"])
    
    # Save the presentation
    output_path = "c:/Users/dario/projects/ai-recommender/docs/Data_Representation_for_Recommendations.pptx"
    prs.save(output_path)
    print(f"Presentation created successfully at: {output_path}")

if __name__ == "__main__":
    create_data_representation_presentation()